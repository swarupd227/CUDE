"""
CUDE Python Parser Worker
Consumes parse jobs from Redis, processes files with Python libraries,
and writes results back to Redis for the Node.js backend to pick up.

Communication protocol:
  - Node.js LPUSH job JSON to 'cude:python:jobs'
  - Worker BRPOP from 'cude:python:jobs', processes file,
  - Worker LPUSH result JSON to 'cude:python:results:{jobId}'
"""

import json
import os
import sys
import time
import traceback
import redis

REDIS_URL = os.environ.get('REDIS_URL', 'redis://localhost:6379')
JOBS_QUEUE = 'cude:python:jobs'
HEARTBEAT_KEY = 'cude:python-worker:heartbeat'

def get_redis():
    return redis.Redis.from_url(REDIS_URL, decode_responses=True)

def parse_pdf(file_path):
    """Enhanced PDF parsing with pdfplumber — better table extraction"""
    import pdfplumber
    result = {
        'parser': 'pdfplumber (Python)',
        'tables': [],
        'text': '',
        'page_count': 0,
    }
    with pdfplumber.open(file_path) as pdf:
        result['page_count'] = len(pdf.pages)
        all_text = []
        for i, page in enumerate(pdf.pages[:20]):  # Limit to 20 pages
            text = page.extract_text() or ''
            all_text.append(text)
            # Extract tables
            tables = page.extract_tables()
            for t in tables:
                if t and len(t) > 1:
                    result['tables'].append({
                        'page': i + 1,
                        'rows': len(t),
                        'cols': len(t[0]) if t[0] else 0,
                        'header': t[0] if t[0] else [],
                        'preview': t[:3],
                    })
        result['text'] = '\n'.join(all_text)[:5000]
    result['table_count'] = len(result['tables'])
    return result

def parse_pptx(file_path):
    """Enhanced PPTX parsing with python-pptx — slide content + notes"""
    from pptx import Presentation
    prs = Presentation(file_path)
    result = {
        'parser': 'python-pptx (Python)',
        'slide_count': len(prs.slides),
        'slides': [],
        'text': '',
        'notes_text': '',
    }
    all_text = []
    notes_text = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)
        # Speaker notes
        note = ''
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                notes_text.append(f'[Slide {i+1} Notes] {note}')

        result['slides'].append({
            'number': i + 1,
            'text': ' '.join(slide_text),
            'has_notes': bool(note),
        })
        all_text.extend(slide_text)

    result['text'] = '\n'.join(all_text)[:5000]
    result['notes_text'] = '\n'.join(notes_text)[:2000]
    result['has_speaker_notes'] = len(notes_text) > 0
    return result

def process_job(job_data):
    """Process a single parse job"""
    job_id = job_data.get('jobId', 'unknown')
    file_path = job_data.get('filePath', '')
    parse_type = job_data.get('parseType', '')

    if not os.path.exists(file_path):
        return {'error': f'File not found: {file_path}', 'jobId': job_id}

    start = time.time()

    try:
        if parse_type == 'pdf':
            result = parse_pdf(file_path)
        elif parse_type == 'pptx':
            result = parse_pptx(file_path)
        else:
            return {'error': f'Unknown parse type: {parse_type}', 'jobId': job_id}

        result['jobId'] = job_id
        result['duration_ms'] = int((time.time() - start) * 1000)
        result['success'] = True
        return result
    except Exception as e:
        return {'error': str(e), 'traceback': traceback.format_exc(), 'jobId': job_id, 'success': False}

def main():
    print('🐍 CUDE Python Parser Worker starting...')
    r = get_redis()

    # Set heartbeat
    r.set(HEARTBEAT_KEY, 'alive', ex=60)
    print(f'📡 Connected to Redis. Listening on queue: {JOBS_QUEUE}')

    while True:
        try:
            # Update heartbeat
            r.set(HEARTBEAT_KEY, 'alive', ex=60)

            # Block-wait for a job (timeout 30s to refresh heartbeat)
            job = r.brpop(JOBS_QUEUE, timeout=30)
            if not job:
                continue

            _, job_json = job
            job_data = json.loads(job_json)
            job_id = job_data.get('jobId', 'unknown')

            print(f'📄 Processing job {job_id}: {job_data.get("parseType")} — {job_data.get("filePath", "")[:60]}')

            result = process_job(job_data)

            # Push result back to Redis
            result_key = f'cude:python:results:{job_id}'
            r.lpush(result_key, json.dumps(result))
            r.expire(result_key, 300)  # Result expires after 5 minutes

            status = '✅' if result.get('success') else '❌'
            print(f'  {status} Job {job_id} done in {result.get("duration_ms", 0)}ms')

        except redis.ConnectionError:
            print('⚠️  Redis connection lost. Retrying in 5s...')
            time.sleep(5)
        except Exception as e:
            print(f'❌ Worker error: {e}')
            time.sleep(1)

if __name__ == '__main__':
    main()
